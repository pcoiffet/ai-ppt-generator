"""
JSON to PowerPoint converter.
Uses centralized Pydantic models from schemas.py
"""
import io
import os
import logging

import requests

logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml import parse_xml

from schemas import (
    PresentationInput,
    SlideContent,
    TextFormatting,
    TextRun,
    BulletPoint,
    TableData,
    ChartData,
    ImageData,
    PPTGenerationError,
)


# === CONFIGURATION ===

# Unsplash API key (free at unsplash.com/developers)
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY", "")

# Fallback image when Unsplash is not available
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FALLBACK_IMAGE = os.path.join(BASE_DIR, "images", "placeholder.jpg")

# Layout name mapping (must match your PowerPoint template)
LAYOUTS = {
    "title_slide": "Title Slide",
    "content_only": "Content Only",
    "image_right": "Image Right",
    "image_left": "Image Left",
    "image_full": "Image Full",
    "table": "Table",
    "chart": "Chart",
    "chart_with_text": "Chart with Text",
    "two_columns": "Two Columns",
}

# Placeholder type constants
PH_TITLE, PH_BODY, PH_OBJECT, PH_CHART, PH_TABLE, PH_PICTURE = 1, 2, 7, 8, 12, 18


# === UTILITIES ===

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def apply_formatting(run, fmt: TextFormatting):
    """Apply text formatting to a run."""
    if fmt.bold: run.font.bold = True
    if fmt.italic: run.font.italic = True
    if fmt.color: run.font.color.rgb = hex_to_rgb(fmt.color)
    if fmt.size: run.font.size = Pt(fmt.size)


def get_layout(prs: Presentation, name: str):
    """Get a slide layout by name from the presentation."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise PPTGenerationError(f"Layout '{name}' not found", 
                             {"available": [l.name for l in prs.slide_layouts]})


def detect_layout(slide: SlideContent) -> str:
    """Auto-detect the best layout for a slide based on its content."""
    # Priority to complex content (Chart/Table)
    if slide.chart: return LAYOUTS["chart"]
    if slide.table: return LAYOUTS["table"]
    
    # Image handling
    if slide.image:
        return LAYOUTS[f"image_{slide.image.position}"]

    # Layout requested by JSON
    if slide.layout and slide.layout in LAYOUTS:
        return LAYOUTS[slide.layout]

    # Fallback
    return LAYOUTS["content_only"]


# === PLACEHOLDER FILLING ===

def fill_text(placeholder, content, bullet_points=None):
    """Fill a text placeholder with content and/or bullet points."""
    tf = placeholder.text_frame
    tf.clear()
    
    if content:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        if isinstance(content, list):
            for run_data in content:
                if isinstance(run_data, dict): run_data = TextRun(**run_data)
                run = p.add_run()
                run.text = run_data.text
                if run_data.formatting: apply_formatting(run, run_data.formatting)
                if run_data.hyperlink:
                    run.hyperlink.address = run_data.hyperlink
                    run.font.color.rgb = RGBColor(0, 0, 255)
                    run.font.underline = True
        else:
            p.text = str(content)
    
    if bullet_points:
        for point in bullet_points:
            bp = BulletPoint(**point) if isinstance(point, dict) else (point if isinstance(point, BulletPoint) else BulletPoint(text=str(point)))
            p = tf.add_paragraph()
            p.text = bp.text
            p.level = bp.level
            pPr = p._pPr
            pPr.append(parse_xml('<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="arabicPeriod"/>'))
            if bp.formatting:
                for run in p.runs: apply_formatting(run, bp.formatting)


def fill_table(placeholder, data: TableData):
    """Fill a table placeholder with data."""
    gf = placeholder.insert_table(len(data.rows) + 1, len(data.headers))
    table = gf.table
    
    for i, header in enumerate(data.headers):
        cell = table.cell(0, i)
        cell.text = str(header)
        if data.style == 'header_colored':
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
    
    for row_idx, row in enumerate(data.rows):
        for col_idx, val in enumerate(row):
            table.cell(row_idx + 1, col_idx).text = str(val)


def fill_chart(placeholder, data: ChartData):
    """Fill a chart placeholder with data."""
    chart_data = CategoryChartData()
    chart_data.categories = data.categories
    for s in data.series:
        chart_data.add_series(s.name, s.data)
    
    chart_types = {'column': XL_CHART_TYPE.COLUMN_CLUSTERED, 'line': XL_CHART_TYPE.LINE, 'pie': XL_CHART_TYPE.PIE}
    placeholder.insert_chart(chart_types.get(data.type, XL_CHART_TYPE.COLUMN_CLUSTERED), chart_data)


def fetch_image_from_unsplash(query: str) -> io.BytesIO | None:
    """Download an image from Unsplash based on a search query."""
    if not UNSPLASH_ACCESS_KEY:
        logger.debug("UNSPLASH_ACCESS_KEY not configured, using fallback")
        return None
    
    try:
        url = "https://api.unsplash.com/search/photos"
        params = {"query": query, "per_page": 1, "orientation": "landscape"}
        headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
        
        response = requests.get(url, params=params, headers=headers, timeout=10)
        response.raise_for_status()
        
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            logger.info(f"Unsplash image found: {image_url[:50]}...")
            img_response = requests.get(image_url, timeout=15)
            img_response.raise_for_status()
            return io.BytesIO(img_response.content)
        else:
            logger.warning(f"No Unsplash image for: {query}")
    except Exception as e:
        logger.warning(f"Unsplash error: {e}")
    
    return None


def fill_image(placeholder, data: ImageData):
    """Insert an image from Unsplash, local path, or fallback placeholder."""
    
    # 1. Try Unsplash
    image_stream = fetch_image_from_unsplash(data.path)
    if image_stream:
        try:
            placeholder.insert_picture(image_stream)
            return
        except Exception as e:
            logger.warning(f"Unsplash image insertion error: {e}")
    
    # 2. Fallback: local image if exists
    if os.path.exists(data.path):
        try:
            placeholder.insert_picture(data.path)
            return
        except Exception as e:
            logger.warning(f"Local image error: {e}")
    
    # 3. Fallback: placeholder image
    if os.path.exists(FALLBACK_IMAGE):
        try:
            placeholder.insert_picture(FALLBACK_IMAGE)
            logger.debug(f"Placeholder image used for: {data.path}")
            return
        except Exception as e:
            logger.warning(f"Placeholder image error: {e}")
    
    # 4. Last fallback: text
    if placeholder.has_text_frame:
        placeholder.text = f"[IMAGE: {data.path}]"


# === MAIN GENERATOR ===

def generate_presentation_stream(json_data: dict, template_stream: io.BytesIO) -> io.BytesIO:
    """Generate a PPTX from JSON data and a template."""
    
    # Validation
    try:
        data = PresentationInput(**json_data)
    except Exception as e:
        raise PPTGenerationError("Validation error", {"details": str(e)})
    
    # Load template
    template_stream.seek(0)
    prs = Presentation(template_stream)
    
    # Clear existing slides
    while prs.slides:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # Metadata
    prs.core_properties.title = data.title
    if data.author: prs.core_properties.author = data.author
    if data.subject: prs.core_properties.subject = data.subject
    
    # Title slide
    slide = prs.slides.add_slide(get_layout(prs, LAYOUTS["title_slide"]))
    body_phs = [ph for ph in slide.placeholders if ph.placeholder_format.type == PH_BODY]
    if body_phs: body_phs[0].text = data.title
    if len(body_phs) > 1 and data.subtitle: body_phs[1].text = data.subtitle
    
    # Content slides
    for slide_data in data.slides:
        layout = get_layout(prs, detect_layout(slide_data))
        slide = prs.slides.add_slide(layout)
        
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type
            
            if ph_type == PH_TITLE:
                ph.text = slide_data.title
            elif ph_type == PH_PICTURE and slide_data.image:
                fill_image(ph, slide_data.image)
            elif ph_type == PH_TABLE and slide_data.table:
                fill_table(ph, slide_data.table)
            elif ph_type == PH_CHART and slide_data.chart:
                fill_chart(ph, slide_data.chart)
            elif ph_type in (PH_BODY, PH_OBJECT) and (slide_data.content or slide_data.bullet_points):
                body_placeholders = [p for p in slide.placeholders if p.placeholder_format.type in (PH_BODY, PH_OBJECT)]
                if body_placeholders and ph.placeholder_format.idx == body_placeholders[0].placeholder_format.idx:
                     fill_text(ph, slide_data.content, slide_data.bullet_points)
    
    # Export
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
